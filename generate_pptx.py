from pptx import Presentation
from pptx.util import Inches, Pt

# Load the markdown presentation
with open('docs/presentacion_tfg.md', 'r', encoding='utf-8') as f:
    content = f.read()

# Parse the markdown into slides
slides_content = []
current_slide = None
for line in content.split('\n'):
    if line.startswith('# '):
        if current_slide is not None:
            slides_content.append(current_slide)
        current_slide = {'title': line[2:], 'body': []}
    elif line.startswith('### '):
        if current_slide is not None:
            current_slide['body'].append(line)
    elif line.strip() == '---' or line.startswith('marp:') or line.startswith('theme:') or line.startswith('paginate:') or line.startswith('header:') or line.startswith('footer:') or line.startswith('style:') or line.startswith('  section ') or line.startswith('  }'):
        continue
    else:
        if current_slide is not None and line.strip() != '':
            current_slide['body'].append(line)

if current_slide is not None:
    slides_content.append(current_slide)

prs = Presentation('plantilla.pptx')

# Modify the title slide
title_slide = prs.slides[0]
for shape in title_slide.shapes:
    if hasattr(shape, "text") and ("Portal Web" in shape.text or "Gestión de" in shape.text):
        shape.text = "Infraestructura Cloud Automatizada y Segura (DevSecOps)"
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(28)

# Safely delete other slides
xml_slides = prs.slides._sldIdLst
slides_to_del = list(xml_slides)[1:]
for sld in slides_to_del:
    xml_slides.remove(sld)

# Extract layout 0 (the only layout)
content_layout = prs.slide_layouts[0]

# Generate other slides
for slide_data in slides_content[1:]:
    slide = prs.slides.add_slide(content_layout)
    
    # Add title text box: 1 inch from left, 0.5 inch from top, 8 inch wide, 1 inch high
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = slide_data['title']
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    
    # Add body text box: 1 inch from left, 2 inches from top, 8 inch wide, 4.5 inch high
    body_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    body_tf.text = '\n'.join(slide_data['body'])
    
    for paragraph in body_tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(20)

prs.save('docs/TFG_Presentacion_Final.pptx')
print("Successfully generated docs/TFG_Presentacion_Final.pptx")
