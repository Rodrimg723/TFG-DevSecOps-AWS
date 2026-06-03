from pptx import Presentation

prs = Presentation('plantilla.pptx')

with open('slides_analysis.txt', 'w', encoding='utf-8') as f:
    f.write(f"Total Slides: {len(prs.slides)}\n")
    for i, slide in enumerate(prs.slides):
        f.write(f"--- Slide {i} ---\n")
        f.write(f"  Shapes:\n")
        for j, shape in enumerate(slide.shapes):
            shape_type = type(shape).__name__
            text = shape.text[:50].replace('\n', ' ') if hasattr(shape, "text") else ""
            f.write(f"    Shape {j} ({shape_type}): {text}\n")
